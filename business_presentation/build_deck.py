"""
LB Marketing — Client Pitch Deck Generator
============================================
Usage:
    python build_deck.py                        # uses client_config.yaml
    python build_deck.py my_other_config.yaml   # uses a custom config file

All editable fields live in client_config.yaml.
"""

import sys
import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ══════════════════════════════════════════════════════════════════════════════
# Load Config
# ══════════════════════════════════════════════════════════════════════════════

def load_config(path="client_config.yaml"):
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


# ══════════════════════════════════════════════════════════════════════════════
# Color Palette
# ══════════════════════════════════════════════════════════════════════════════

NAVY      = RGBColor(0x1E, 0x27, 0x61)
TEAL      = RGBColor(0x02, 0x80, 0x90)
ICE       = RGBColor(0xCA, 0xDC, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF7, 0xFF)
CHARCOAL  = RGBColor(0x2E, 0x2E, 0x3A)
MID_GRAY  = RGBColor(0x64, 0x74, 0x8B)
AMBER     = RGBColor(0xF5, 0xA6, 0x23)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
DARK_BLUE = RGBColor(0x28, 0x3A, 0x5E)
RED       = RGBColor(0xDC, 0x26, 0x26)
ROW_ALT   = RGBColor(0xF0, 0xF7, 0xFF)
LB_COL    = RGBColor(0xE0, 0xF4, 0xF6)

W = Inches(10)
H = Inches(5.625)


# ══════════════════════════════════════════════════════════════════════════════
# Drawing Helpers
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=14, bold=False, italic=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox


def nav_bar(slide, label, agency_name):
    add_rect(slide, 0, 0, W, Inches(0.52), NAVY)
    add_text(slide, agency_name, Inches(0.3), Inches(0.08), Inches(3), Inches(0.38),
             size=13, bold=True, color=WHITE)
    if label:
        add_text(slide, label, Inches(0), Inches(0.08), W - Inches(0.3), Inches(0.38),
                 size=11, color=ICE, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), NAVY)
    add_text(slide, text, Inches(0.3), H - Inches(0.27), W - Inches(0.6), Inches(0.26),
             size=8, color=ICE)


def stat_block(slide, x, y, w, h, number, label, bg=NAVY):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, number, x, y + Inches(0.08), w, Inches(0.5),
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.52), w, Inches(0.35),
             size=9, color=ICE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide Builders
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs, cfg, BLANK):
    ag = cfg["agency"]
    cl = cfg["client"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, NAVY)
    add_rect(sl, 0, Inches(2.8), W, Inches(0.08), TEAL)
    add_text(sl, ag["name"].upper(), Inches(0.6), Inches(0.8), Inches(8.8), Inches(0.8),
             size=42, bold=True, color=WHITE, font="Calibri")
    add_text(sl, ag["tagline"], Inches(0.6), Inches(1.55), Inches(8.8), Inches(0.5),
             size=18, color=ICE)
    add_rect(sl, Inches(0.6), Inches(3.0), Inches(5.5), Inches(1.1), TEAL)
    add_text(sl, "Presented to:", Inches(0.75), Inches(3.08), Inches(5), Inches(0.3),
             size=10, color=ICE)
    add_text(sl, cl["name"], Inches(0.75), Inches(3.32), Inches(5), Inches(0.45),
             size=20, bold=True, color=WHITE)
    add_text(sl, f"{cl['city_state']}  •  {cl['date']}",
             Inches(0.75), Inches(3.72), Inches(5), Inches(0.3), size=10, color=ICE)
    for i, (num, lbl) in enumerate([
        ("Month 1", "No Lock-In Demo"),
        ("30 Days", "To First Results"),
        ("5-10",    "New Reviews Target"),
    ]):
        stat_block(sl, Inches(6.8), Inches(3.0) + i * Inches(0.68),
                   Inches(2.8), Inches(0.62), num, lbl, CHARCOAL)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_02_about(prs, cfg, BLANK):
    ag = cfg["agency"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "About Us", ag["name"])
    add_text(sl, "Who We Are", Inches(0.4), Inches(0.65), Inches(6), Inches(0.55),
             size=28, bold=True, color=NAVY)
    add_text(sl,
        f"{ag['name']} is a {ag['location']}-based local digital marketing agency built "
        "specifically for small and medium businesses in three verticals: Salons & Spas, "
        "Legal & Professional Services, and Home Services. We combine hands-on "
        "agency work with a proprietary software platform to deliver measurable results -- fast.",
        Inches(0.4), Inches(1.25), Inches(5.6), Inches(1.4),
        size=12, color=CHARCOAL, wrap=True)
    pillars = [
        ("Results First",        "Tangible outcomes by end of Week 4 -- before you commit to anything more."),
        ("No Lock-In",           "Month-to-month structure. Earn your trust every month."),
        ("Vertical Expertise",   "Deep focus on 3 industries means faster wins, not generic strategy."),
        ("Proprietary Platform", "Built-in tools for audit reports, review campaigns & month-end reporting."),
    ]
    for i, (title, body) in enumerate(pillars):
        col, row = i % 2, i // 2
        cx = Inches(0.4) + col * Inches(4.85)
        cy = Inches(2.75) + row * Inches(1.1)
        add_rect(sl, cx, cy, Inches(4.55), Inches(0.95), WHITE)
        add_rect(sl, cx, cy, Inches(0.07), Inches(0.95), TEAL)
        add_text(sl, title, cx + Inches(0.18), cy + Inches(0.08), Inches(4.2), Inches(0.28),
                 size=12, bold=True, color=NAVY)
        add_text(sl, body, cx + Inches(0.18), cy + Inches(0.35), Inches(4.2), Inches(0.55),
                 size=10, color=CHARCOAL)
    add_rect(sl, Inches(6.2), Inches(0.6), Inches(3.4), Inches(1.9), NAVY)
    add_text(sl, "Our Approach", Inches(6.35), Inches(0.68), Inches(3.1), Inches(0.35),
             size=12, bold=True, color=WHITE)
    for i, pt in enumerate([
        "Audit -> Fix -> Measure",
        "Leads with leading indicators only",
        "Specific gaps, not generic pitches",
        "Platform-powered delivery",
    ]):
        add_text(sl, f"  checkmark  {pt}", Inches(6.35), Inches(1.05) + i * Inches(0.33),
                 Inches(3.1), Inches(0.3), size=10, color=ICE)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_03_analysis(prs, cfg, BLANK):
    ag = cfg["agency"]
    cl = cfg["client"]
    cs = cfg["current_state"]
    op = cfg["opportunity"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Your Business Analysis", ag["name"])
    add_text(sl, f"What We Found: {cl['name']}",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_text(sl, "A snapshot of your current online presence -- the gaps and the opportunities.",
             Inches(0.4), Inches(1.12), Inches(9.2), Inches(0.3),
             size=11, color=MID_GRAY, italic=True)
    add_rect(sl, Inches(0.4), Inches(1.5), Inches(4.5), Inches(3.5), WHITE)
    add_rect(sl, Inches(0.4), Inches(1.5), Inches(0.07), Inches(3.5), RED)
    add_text(sl, "Current State", Inches(0.57), Inches(1.58), Inches(4.2), Inches(0.35),
             size=13, bold=True, color=RED)
    current_rows = [
        ("Google Business Profile", cs.get("gbp_status",           "--")),
        ("Star Rating",             cs.get("star_rating",          "--")),
        ("Review Count",            cs.get("review_count",         "--")),
        ("Map Pack Rank",           cs.get("map_pack_rank",        "--")),
        ("Website",                 cs.get("website",              "--")),
        ("Citation Consistency",    cs.get("citation_consistency", "--")),
        ("Last Review Response",    cs.get("last_review_response", "--")),
    ]
    for i, (label, val) in enumerate(current_rows):
        y = Inches(2.0) + i * Inches(0.38)
        add_text(sl, label, Inches(0.6),  y, Inches(2.2), Inches(0.35), size=10, color=CHARCOAL, bold=True)
        add_text(sl, val,   Inches(2.85), y, Inches(2.0), Inches(0.35), size=10, color=MID_GRAY)
    add_rect(sl, Inches(5.1), Inches(1.5), Inches(4.5), Inches(3.5), WHITE)
    add_rect(sl, Inches(5.1), Inches(1.5), Inches(0.07), Inches(3.5), TEAL)
    add_text(sl, "Opportunity", Inches(5.27), Inches(1.58), Inches(4.2), Inches(0.35),
             size=13, bold=True, color=TEAL)
    opp_rows = [
        ("Map Pack Opportunity",     op.get("map_pack_opportunity", "--")),
        ("Review Gap vs. Top 3",     op.get("review_gap",           "--")),
        ("Keyword Gap",              op.get("keyword_gap",          "--")),
        ("Profile Completion Score", op.get("profile_completion",   "--")),
        ("Review Velocity",          op.get("review_velocity",      "--")),
        ("Unanswered Reviews",       op.get("unanswered_reviews",   "--")),
        ("Quick Win Potential",      op.get("quick_win_potential",  "--")),
    ]
    for i, (label, val) in enumerate(opp_rows):
        y = Inches(2.0) + i * Inches(0.38)
        add_text(sl, label, Inches(5.27), y, Inches(2.0), Inches(0.35), size=10, color=CHARCOAL, bold=True)
        add_text(sl, val,   Inches(7.3),  y, Inches(2.2), Inches(0.35), size=10, color=MID_GRAY)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_04_competitors(prs, cfg, BLANK):
    ag     = cfg["agency"]
    rows   = cfg.get("competitors", [])
    insight = cfg.get("competitor_key_insight", "")
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Local Competitor Analysis", ag["name"])
    add_text(sl, "How You Stack Up Against Local Competitors",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_text(sl, "Data pulled from Google & Yelp for your primary service area.",
             Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.3),
             size=9, color=MID_GRAY, italic=True)
    headers  = ["Business", "Google Rating", "Review Count", "Map Pack?", "GBP Complete?", "Website?"]
    col_keys = ["name", "google_rating", "review_count", "map_pack", "gbp_complete", "website"]
    col_w    = [Inches(2.1), Inches(1.2), Inches(1.2), Inches(1.1), Inches(1.3), Inches(1.0)]
    x_starts = [Inches(0.35)]
    for cw in col_w[:-1]:
        x_starts.append(x_starts[-1] + cw)
    header_y = Inches(1.5)
    for hdr, cw, cx in zip(headers, col_w, x_starts):
        add_rect(sl, cx, header_y, cw, Inches(0.32), NAVY)
        add_text(sl, hdr, cx + Inches(0.05), header_y + Inches(0.04),
                 cw - Inches(0.08), Inches(0.26), size=9, bold=True, color=WHITE)
    for ri, row_data in enumerate(rows[:5]):
        ry = Inches(1.82) + ri * Inches(0.42)
        is_client = (ri == 0)
        row_bg    = TEAL if is_client else (ROW_ALT if ri % 2 == 0 else WHITE)
        row_txt   = WHITE if is_client else CHARCOAL
        for j, (key, cw, cx) in enumerate(zip(col_keys, col_w, x_starts)):
            cell_bg  = TEAL if is_client else (LB_COL if j == 0 else row_bg)
            cell_val = row_data.get(key, "--")
            if is_client and j == 0:
                cell_val = f"{cell_val}  <- You"
            add_rect(sl, cx, ry, cw, Inches(0.4), cell_bg)
            add_text(sl, cell_val, cx + Inches(0.05), ry + Inches(0.07),
                     cw - Inches(0.08), Inches(0.3), size=9,
                     color=row_txt, bold=is_client,
                     align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    add_rect(sl, Inches(0.35), Inches(4.45), Inches(9.3), Inches(0.72), NAVY)
    add_text(sl, "Key Insight:", Inches(0.5), Inches(4.52), Inches(1.5), Inches(0.3),
             size=10, bold=True, color=AMBER)
    add_text(sl, str(insight).strip(), Inches(2.0), Inches(4.52), Inches(7.5), Inches(0.6),
             size=9, color=ICE)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_05_month1(prs, cfg, BLANK):
    ag = cfg["agency"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Month 1 Delivery", ag["name"])
    add_text(sl, "Your First 30 Days -- What We Deliver",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    weeks = [
        ("Week 1", NAVY, "Full Digital Presence Audit", [
            "- Google Business Profile, Yelp, website, reviews, local rankings reviewed",
            "- Written audit report delivered by Day 5 -- prioritised action list included",
        ]),
        ("Week 2", TEAL, "GBP Optimisation + Review Response", [
            "- Every GBP field completed: categories, services, hours, photos, description",
            "- All unanswered reviews responded to professionally (positive & negative)",
        ]),
        ("Week 3", NAVY, "Review Generation Campaign + Content", [
            "- Personalised SMS/email review requests sent to recent customers",
            "- Target: 5-10 new Google reviews in Month 1",
            "- First content published (vertical-specific -- see next slide)",
        ]),
        ("Week 4", TEAL, "Month-End Report", [
            "- Before/after comparison across every metric tracked",
            "- Rankings, review count & rating, profile views -- clearly documented",
        ]),
    ]
    for i, (week, bg, title, bullets) in enumerate(weeks):
        col, row = i % 2, i // 2
        cx = Inches(0.35) + col * Inches(4.85)
        cy = Inches(1.3) + row * Inches(1.95)
        add_rect(sl, cx, cy, Inches(4.55), Inches(1.85), WHITE)
        add_rect(sl, cx, cy, Inches(1.1), Inches(0.38), bg)
        add_text(sl, week, cx + Inches(0.08), cy + Inches(0.06),
                 Inches(1.0), Inches(0.28), size=11, bold=True, color=WHITE)
        add_text(sl, title, cx + Inches(0.12), cy + Inches(0.44),
                 Inches(4.3), Inches(0.3), size=12, bold=True, color=NAVY)
        for bi, bul in enumerate(bullets):
            add_text(sl, bul, cx + Inches(0.12), cy + Inches(0.75) + bi * Inches(0.34),
                     Inches(4.3), Inches(0.32), size=10, color=CHARCOAL)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_06_vertical(prs, cfg, BLANK):
    ag = cfg["agency"]
    v  = cfg["vertical"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Vertical-Specific Deliverables", ag["name"])
    add_text(sl, "Month 1 Extras: Tailored to Your Industry",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_text(sl,
        "The universal Month 1 package plus the actions that move the needle fastest for your specific vertical.",
        Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.3),
        size=10, color=MID_GRAY, italic=True)
    all_verticals = [
        ("salons_spas",        "Salons & Spas",                                     NAVY, [
            "Instagram audit + first 4 real-work posts published",
            "GBP photos refreshed -- minimum 10 high-quality images",
            "Booking link added to Google listing & social profiles",
            "Competitor review gap analysis vs. top 3 nearby salons",
            "Quick-win target: 5+ new Google reviews, profile views up",
        ]),
        ("legal_professional", "Legal & Professional Services",                      TEAL, [
            "Local keyword ranking baseline (top 5 practice-area + city terms)",
            "GBP services section fully built out with practice areas",
            "On-page SEO quick-fixes applied to homepage & top service page",
            "Post-matter review request process set up + first batch sent",
            "Quick-win target: map pack entry for 1+ core keyword by Day 30",
        ]),
        ("home_services",      "Home Services (Electricians, Plumbers, Landscapers)", PURPLE, [
            "Full GBP build: all service types, service area, photos added",
            "Local citation audit -- NAP corrected across directories",
            "Top 5 high-intent keyword ranking baseline established",
            "Review campaign targeting recent job completions",
            "Quick-win target: map pack for 1+ core keyword + first new reviews",
        ]),
    ]
    active = [(key, title, bg, items)
              for key, title, bg, items in all_verticals if v.get(key, False)]
    if not active:
        active = [(key, title, bg, items) for key, title, bg, items in all_verticals]
    count = len(active)
    cw    = Inches(9.3 / count - 0.15)
    gap   = Inches(0.15)
    for i, (key, title, bg, items) in enumerate(active):
        cx = Inches(0.35) + i * (cw + gap)
        cy = Inches(1.5)
        ch = Inches(3.7)
        add_rect(sl, cx, cy, cw, ch, WHITE)
        add_rect(sl, cx, cy, cw, Inches(0.38), bg)
        add_text(sl, title, cx + Inches(0.1), cy + Inches(0.06),
                 cw - Inches(0.15), Inches(0.28), size=10, bold=True, color=WHITE)
        row_h = Inches(3.0 / len(items))
        for bi, item in enumerate(items):
            add_text(sl, f"  {item}", cx + Inches(0.1), cy + Inches(0.5) + bi * row_h,
                     cw - Inches(0.15), row_h, size=9, color=CHARCOAL)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_07_archetype(prs, cfg, BLANK):
    ag = cfg["agency"]
    cl = cfg["client"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Your Starting Point", ag["name"])
    add_text(sl, "We Adapt To Where You Are Today",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_rect(sl, Inches(0.35), Inches(1.45), Inches(4.55), Inches(3.62), NAVY)
    add_text(sl, "Established but Invisible",
             Inches(0.5), Inches(1.55), Inches(4.25), Inches(0.4),
             size=14, bold=True, color=WHITE)
    add_text(sl,
        "You have a solid business but little or no online presence. "
        "Customers find you through referrals, not search.",
        Inches(0.5), Inches(1.95), Inches(4.25), Inches(0.55),
        size=10, color=ICE)
    for i, (wk, desc) in enumerate([
        ("Week 1", "Audit + full GBP build-out from scratch"),
        ("Week 2", "Citations, categories, photos, service listings"),
        ("Week 3", "First review generation campaign to existing customers"),
        ("Week 4", "First reviews in -- profile views increasing"),
    ]):
        y = Inches(2.58) + i * Inches(0.55)
        add_rect(sl, Inches(0.5), y, Inches(0.75), Inches(0.28), TEAL)
        add_text(sl, wk, Inches(0.52), y + Inches(0.03),
                 Inches(0.72), Inches(0.23), size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(1.3), y, Inches(3.45), Inches(0.28), size=9, color=ICE)
    add_rect(sl, Inches(5.1), Inches(1.45), Inches(4.55), Inches(3.62), TEAL)
    add_text(sl, "Present but Underperforming",
             Inches(5.25), Inches(1.55), Inches(4.25), Inches(0.4),
             size=14, bold=True, color=WHITE)
    add_text(sl,
        "You've tried -- a website, a Google listing, maybe social -- "
        "but results are inconsistent and your rating is holding you back.",
        Inches(5.25), Inches(1.95), Inches(4.25), Inches(0.55),
        size=10, color=ICE)
    for i, (wk, desc) in enumerate([
        ("Week 1", "Audit -- emphasis on what's actively hurting you"),
        ("Week 2", "Respond to negative reviews; correct inconsistent data"),
        ("Week 3", "Review velocity campaign -- push rating above 4.0"),
        ("Week 4", "Star rating improved; negative review ratio reduced"),
    ]):
        y = Inches(2.58) + i * Inches(0.55)
        add_rect(sl, Inches(5.25), y, Inches(0.75), Inches(0.28), NAVY)
        add_text(sl, wk, Inches(5.27), y + Inches(0.03),
                 Inches(0.72), Inches(0.23), size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(6.05), y, Inches(3.45), Inches(0.28), size=9, color=ICE)
    ei  = cl.get("archetype_established_invisible",   False)
    pi  = cl.get("archetype_present_underperforming", False)
    mi  = cl.get("archetype_mixed",                   False)
    selector = (f"  {cl['name']} best fits:  "
                f"{'[X]' if ei else '[ ]'} Established but Invisible     "
                f"{'[X]' if pi else '[ ]'} Present but Underperforming     "
                f"{'[X]' if mi else '[ ]'} Mix of both")
    add_rect(sl, Inches(0.35), Inches(5.05), Inches(9.3), Inches(0.22), AMBER)
    add_text(sl, selector, Inches(0.5), Inches(5.06),
             Inches(9.0), Inches(0.20), size=8, bold=True, color=NAVY)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_08_requirements(prs, cfg, BLANK):
    ag = cfg["agency"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "What We Need From You", ag["name"])
    add_text(sl, "Three Things -- That's It",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_text(sl, "We handle everything else. Here's your only to-do list for Month 1.",
             Inches(0.4), Inches(1.1), Inches(6.5), Inches(0.3),
             size=11, color=MID_GRAY, italic=True)
    asks = [
        ("1", "Google Business Profile Access",
         "Owner-level access so we can make changes directly. Takes 2 minutes to grant."),
        ("2", "20-Minute Onboarding Call -- Week 1",
         "We learn your business priorities, confirm your target keywords, and set reporting preferences."),
        ("3", "A Handful of Customer Names / Emails",
         "For the Week 3 review request campaign. We draft every message -- you just provide the list."),
    ]
    for i, (num, title, body) in enumerate(asks):
        cy = Inches(1.6) + i * Inches(1.1)
        add_rect(sl, Inches(0.35), cy, Inches(0.72), Inches(0.72), TEAL)
        add_text(sl, num, Inches(0.35), cy + Inches(0.1),
                 Inches(0.72), Inches(0.52), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(1.15), cy, Inches(8.5), Inches(0.72), WHITE)
        add_rect(sl, Inches(1.15), cy, Inches(0.06), Inches(0.72), TEAL)
        add_text(sl, title, Inches(1.3), cy + Inches(0.06),
                 Inches(8.2), Inches(0.3), size=13, bold=True, color=NAVY)
        add_text(sl, body, Inches(1.3), cy + Inches(0.38),
                 Inches(8.2), Inches(0.3), size=10, color=CHARCOAL)
    add_rect(sl, Inches(0.35), Inches(4.65), Inches(9.3), Inches(0.68), NAVY)
    add_text(sl, "No Lock-In.", Inches(0.55), Inches(4.73), Inches(2.0), Inches(0.35),
             size=14, bold=True, color=AMBER)
    add_text(sl,
        "Month 1 is month-to-month. If we haven't earned your confidence by end of Week 4, "
        "you can walk away -- no penalties, no awkward conversations.",
        Inches(2.55), Inches(4.73), Inches(7.0), Inches(0.52), size=10, color=ICE)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_09_service_comparison(prs, cfg, BLANK):
    ag = cfg["agency"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Why LB Marketing", ag["name"])
    add_text(sl, "How We Compare to Other Marketing Services",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    headers = ["Feature / Provider", ag["name"], "Hibu", "Generic Agency", "DIY / In-House"]
    col_w2  = [Inches(2.3), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)]
    x2 = [Inches(0.35)]
    for cw in col_w2[:-1]:
        x2.append(x2[-1] + cw)
    hy = Inches(1.25)
    for j, (hdr, cw, cx) in enumerate(zip(headers, col_w2, x2)):
        bg = TEAL if j == 1 else NAVY
        add_rect(sl, cx, hy, cw, Inches(0.35), bg)
        add_text(sl, hdr, cx + Inches(0.05), hy + Inches(0.05),
                 cw - Inches(0.08), Inches(0.27),
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    comp_rows = [
        ("Contract / Lock-In",          "None -- month-to-month",  "12-24 month contracts", "Often 6-12 months",  "N/A"),
        ("Local SEO Expertise",          "Deep -- 3 verticals",     "Generalist",             "Varies widely",      "Limited"),
        ("GBP Optimisation",             "Week 2",                  "(slow)",                 "Sometimes",          "Manual"),
        ("Review Generation Campaign",   "Week 3",                  "Limited",                "Add-on cost",        "Time-intensive"),
        ("Month-End Reporting",          "Week 4",                  "Automated only",         "Varies",             "None"),
        ("Results in First 30 Days",     "Guaranteed focus",        "Unlikely",               "Setup phase only",   "Unlikely"),
        ("Proprietary Platform / Tools", "Included",                "No",                     "No",                 "No"),
        ("Trustpilot / Review Score",    "N/A -- new agency",       "~2.5 / 5",               "Varies",             "N/A"),
    ]
    for ri, row in enumerate(comp_rows):
        ry = Inches(1.62) + ri * Inches(0.42)
        bg_row = ROW_ALT if ri % 2 == 0 else WHITE
        for j, (cell, cw, cx) in enumerate(zip(row, col_w2, x2)):
            bg = LB_COL if j == 1 else bg_row
            add_rect(sl, cx, ry, cw, Inches(0.4), bg)
            add_text(sl, cell, cx + Inches(0.05), ry + Inches(0.07),
                     cw - Inches(0.08), Inches(0.3), size=9, color=CHARCOAL,
                     bold=(j == 0), align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_10_pricing(prs, cfg, BLANK):
    ag = cfg["agency"]
    pr = cfg.get("pricing", {})
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "Pricing", ag["name"])
    add_text(sl, "Investment & Service Tiers",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_text(sl, "Month 1 is the demo. Continued tiers reflect the ongoing scope of work as your presence grows.",
             Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.3),
             size=10, color=MID_GRAY, italic=True)
    tiers = [
        ("Month 1 Demo",       pr.get("month1_price", "$[X,XXX]"),    TEAL,   True, [
            "Full Digital Presence Audit (written report)",
            "Google Business Profile Optimisation",
            "Review Response Management",
            "Review Generation Campaign (5-10 new reviews target)",
            "Month-End Before/After Report",
            "Vertical-specific quick-win actions",
            "No contract -- cancel anytime",
        ]),
        ("Growth  (Month 2+)", pr.get("growth_price", "$[X,XXX]/mo"), NAVY,   False, [
            "Everything in Month 1, ongoing",
            "Content programme (blog, social, local posts)",
            "Ongoing SEO -- keyword tracking & improvements",
            "Referral & local partnership outreach (Month 3)",
            "Monthly reporting cadence",
            "[WIP] Tier in development -- details to follow",
        ]),
        ("Scale  (Month 4+)",  pr.get("scale_price",  "$[X,XXX]/mo"), PURPLE, False, [
            "Everything in Growth, plus:",
            "Paid campaign testing (where appropriate)",
            "Full-funnel visibility dashboard",
            "Priority support & quarterly strategy review",
            "[WIP] Tier in development -- details to follow",
        ]),
    ]
    for i, (name, price, bg, featured, items) in enumerate(tiers):
        cx = Inches(0.35) + i * Inches(3.2)
        cy = Inches(1.45)
        cw = Inches(3.0)
        add_rect(sl, cx, cy, cw, Inches(3.65), WHITE)
        add_rect(sl, cx, cy, cw, Inches(0.55), bg)
        add_text(sl, name, cx + Inches(0.1), cy + Inches(0.06),
                 cw - Inches(0.15), Inches(0.28), size=11, bold=True, color=WHITE)
        add_text(sl, price, cx + Inches(0.1), cy + Inches(0.32),
                 cw - Inches(0.15), Inches(0.22), size=10, color=ICE)
        if featured:
            add_rect(sl, cx + cw - Inches(1.15), cy + Inches(0.06),
                     Inches(1.1), Inches(0.22), AMBER)
            add_text(sl, "START HERE",
                     cx + cw - Inches(1.15), cy + Inches(0.07),
                     Inches(1.1), Inches(0.2),
                     size=7, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        for bi, item in enumerate(items):
            is_wip = "[WIP]" in item
            label  = item.replace("[WIP] ", "") if is_wip else f"  {item}"
            add_text(sl, label,
                     cx + Inches(0.1), cy + Inches(0.65) + bi * Inches(0.44),
                     cw - Inches(0.15), Inches(0.42),
                     size=9, color=MID_GRAY if is_wip else CHARCOAL, italic=is_wip)
    add_text(sl,
        "All tiers are month-to-month. Growth & Scale tiers are actively being refined -- "
        "pricing and inclusions subject to change.",
        Inches(0.35), Inches(5.2), Inches(9.3), Inches(0.2),
        size=8, color=MID_GRAY, italic=True)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_11_roadmap(prs, cfg, BLANK):
    ag = cfg["agency"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, OFF_WHITE)
    nav_bar(sl, "The Roadmap", ag["name"])
    add_text(sl, "Month 1 Is Just the Beginning",
             Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.5),
             size=24, bold=True, color=NAVY)
    add_text(sl, "Leading indicators (Month 1) unlock lagging outcomes (Months 2-4+). Here's the full picture.",
             Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.3),
             size=10, color=MID_GRAY, italic=True)
    timeline = [
        ("Month 1 (Demo)", TEAL, "Fast Foundation", [
            "GBP fully optimised", "Review velocity started",
            "Ranking baseline set", "Before/after report",
        ]),
        ("Month 2", NAVY, "Content & SEO Engine", [
            "Regular posting & blog", "Ongoing local SEO",
            "Organic traffic grows", "Monthly reporting",
        ]),
        ("Month 3", TEAL, "Referral & Partnerships", [
            "Referral programme", "Local biz partnerships",
            "Word-of-mouth amplified", "Review momentum",
        ]),
        ("Month 4+", NAVY, "Full-Funnel Visibility", [
            "Paid campaign testing", "Full-funnel dashboard",
            "Search to enquiry to client", "Quarterly strategy",
        ]),
    ]
    for i, (period, bg, subtitle, pts) in enumerate(timeline):
        cx = Inches(0.35) + i * Inches(2.4)
        cy = Inches(1.45)
        cw = Inches(2.2)
        add_rect(sl, cx, cy, cw, Inches(0.55), bg)
        add_text(sl, period, cx + Inches(0.08), cy + Inches(0.05),
                 cw - Inches(0.1), Inches(0.48), size=11, bold=True, color=WHITE)
        add_rect(sl, cx, cy + Inches(0.55), cw, Inches(3.35), WHITE)
        add_text(sl, subtitle, cx + Inches(0.1), cy + Inches(0.62),
                 cw - Inches(0.15), Inches(0.3), size=10, bold=True, color=NAVY)
        for bi, pt in enumerate(pts):
            add_text(sl, f"  {pt}",
                     cx + Inches(0.1), cy + Inches(0.98) + bi * Inches(0.58),
                     cw - Inches(0.15), Inches(0.5), size=9, color=CHARCOAL)
        if i < 3:
            add_rect(sl, cx + cw, cy + Inches(0.2), Inches(0.2), Inches(0.22), OFF_WHITE)
            add_text(sl, ">", cx + cw, cy + Inches(0.2),
                     Inches(0.2), Inches(0.22), size=12, color=TEAL, align=PP_ALIGN.CENTER)
    footer(sl, f"Confidential  -  {ag['name']}  -  {ag['location']}")


def slide_12_close(prs, cfg, BLANK):
    ag = cfg["agency"]
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, NAVY)
    add_rect(sl, 0, Inches(2.5), W, Inches(0.08), TEAL)
    add_text(sl, "Ready to See Results in 30 Days?",
             Inches(0.6), Inches(0.75), Inches(8.8), Inches(0.7),
             size=32, bold=True, color=WHITE)
    add_text(sl,
        "Let's book your onboarding call. We'll confirm access, align on priorities, "
        "and start the audit immediately. No paperwork, no commitment beyond Month 1.",
        Inches(0.6), Inches(1.45), Inches(7.5), Inches(0.9),
        size=13, color=ICE, wrap=True)
    add_rect(sl, Inches(0.6), Inches(2.7), Inches(4.5), Inches(1.6), TEAL)
    add_text(sl, "Get In Touch", Inches(0.75), Inches(2.78),
             Inches(4.2), Inches(0.35), size=13, bold=True, color=WHITE)
    for i, line in enumerate([
        ag["presenter_name"],
        ag["presenter_email"],
        ag["presenter_phone"],
        ag["website"],
    ]):
        add_text(sl, line, Inches(0.75), Inches(3.14) + i * Inches(0.3),
                 Inches(4.2), Inches(0.28), size=11, color=ICE)
    add_rect(sl, Inches(5.4), Inches(2.7), Inches(4.2), Inches(1.6), DARK_BLUE)
    add_text(sl, "Next Steps", Inches(5.55), Inches(2.78),
             Inches(3.9), Inches(0.35), size=13, bold=True, color=WHITE)
    for i, step in enumerate([
        "1.  Sign off on access (GBP)",
        "2.  Book 20-min onboarding call",
        "3.  We start Day 1",
    ]):
        add_text(sl, step, Inches(5.55), Inches(3.14) + i * Inches(0.3),
                 Inches(3.9), Inches(0.28), size=11, color=ICE)
    add_text(sl, "No lock-in  -  Month-to-month  -  Results by Week 4",
             Inches(0), Inches(4.5), W, Inches(0.4),
             size=11, color=AMBER, align=PP_ALIGN.CENTER, bold=True)
    footer(sl, f"{ag['name']}  -  {ag['location']}  -  Confidential")


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

def build(config_path="client_config.yaml"):
    cfg = load_config(config_path)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    BLANK = prs.slide_layouts[6]

    slide_01_cover(prs, cfg, BLANK)
    slide_02_about(prs, cfg, BLANK)
    slide_03_analysis(prs, cfg, BLANK)
    slide_04_competitors(prs, cfg, BLANK)
    slide_05_month1(prs, cfg, BLANK)
    slide_06_vertical(prs, cfg, BLANK)
    slide_07_archetype(prs, cfg, BLANK)
    slide_08_requirements(prs, cfg, BLANK)
    slide_09_service_comparison(prs, cfg, BLANK)
    slide_10_pricing(prs, cfg, BLANK)
    slide_11_roadmap(prs, cfg, BLANK)
    slide_12_close(prs, cfg, BLANK)

    client_name  = cfg["client"]["name"]
    safe_name    = "".join(c if c.isalnum() or c in "_- " else "" for c in client_name).strip().replace(" ", "_")
    out_template = cfg.get("output_filename", "LB_Marketing_Pitch_[ClientName].pptx")
    out_path     = out_template.replace("[ClientName]", safe_name)

    prs.save(out_path)
    print(f"  Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    config_file = sys.argv[1] if len(sys.argv) > 1 else "client_config.yaml"
    if not os.path.exists(config_file):
        print(f"  Config file not found: {config_file}")
        sys.exit(1)
    build(config_file)
